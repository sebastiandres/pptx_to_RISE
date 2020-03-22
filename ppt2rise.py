# Ressources
# 1. https://python-pptx.readthedocs.io/en/latest/
# 2. https://stackoverflow.com/questions/32908639/open-pil-image-from-byte-file
# 2. Shapes: https://python-pptx.readthedocs.io/en/latest/user/autoshapes.html

import sys
import os
import io
import json
from PIL import Image
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches

# SOME PARAMETERS
DEBUG = False

def create_notebook(cells):
    """
    Creates a jupyter notebook given a list of cells
    """
    notebook_dict = {}
    notebook_dict["cells"] = cells
    notebook_dict["metadata"] = {
                                "celltoolbar": "Slideshow",
                                "kernelspec": {
                                    "display_name": "Python 3",
                                    "language": "python",
                                    "name": "python3"
                                    },
                                "language_info": {
                                    "codemirror_mode": {"name": "ipython","version": 3},
                                    "file_extension": ".py",
                                    "mimetype": "text/x-python",
                                    "name": "python",
                                    "nbconvert_exporter": "python",
                                    "pygments_lexer": "ipython3",
                                    "version": "3.7.4"
                                    }
                                }
    notebook_dict["nbformat"] = 4
    notebook_dict["nbformat_minor"] = 2
    return notebook_dict

def save_image(filepath, blob):
    """
    """    
    image = Image.open(io.BytesIO(blob))
    image.save(filepath)
    return

def save(notebook_dict, ipynb_filepath):
    with open(ipynb_filepath, 'w') as fp:
        json.dump(notebook_dict, fp)
    return

def create_cell(cell_type="markdown", 
                content="This is markdown!", 
                slide_type="slide"):
    """
    This function creates a cell
    """
    cell = {
                "cell_type": cell_type,
                "metadata": {
                    "slideshow": {
                    "slide_type": slide_type
                    }
                            },
                "source": content
            }
    if cell_type=="code":
        cell["execution_count"] = None
        cell["outputs"] = []
    return cell    

def config_cell():
    code = """#!/usr/bin/env python3
from traitlets.config.manager import BaseJSONConfigManager
from pathlib import Path
path = Path.home() / ".jupyter" / "nbconfig"
cm = BaseJSONConfigManager(config_dir=str(path))
cm.update(
    "rise",
    {
        "theme": "none", # sky, ...
        "transition": "none", #
    }
)"""
    return create_cell(cell_type="code", 
                       content=code, slide_type="skip")

def create_image_markdown(shape, output_folder):
    """
    Creates an image and the markdown content.
    """
    blob = shape.image.blob
    content_type = shape.image.content_type.split("/")[-1]
    key = datetime.utcnow().strftime('%Y_%m_%d_%H_%M_%S_%f')
    image_filename = "{}/{}.{}".format("images", key, content_type)
    image_filepath = "{}/{}".format(output_folder, image_filename)
    save_image(image_filepath, blob)
    my_text = '\n !["{}"]({}) \n'.format("Image", image_filename) # Not filepath, just the name because same folder
    return my_text

def get_markdown(shape, output_folder, preppend=""):
    """
    """
    if "text" in dir(shape):
        my_text =  "".join([preppend+_+"\n" for _ in shape.text.split("\n")])
    elif "image" in dir(shape):
        my_text = create_image_markdown(shape, output_folder)
    else:
        my_text = ""
    return my_text

def mkdir(my_dir):
    try:
        os.mkdir(my_dir)
    except:
        pass

def get_type(shape):
    if "text" in dir(shape):
        return "text"
    elif "image" in dir(shape):
        return "image"
    else:
        return "unknown"

def is_left_column(shape, slide_width):
    return shape.left.inches<0.5*slide_width

def is_right_column(shape, slide_width):
    return shape.left.inches>0.5*slide_width

def ppt2rise(input_filepath, output_filepath, debug=DEBUG):

    # Check sanity of arguments
    output_folder = os.path.dirname(output_filepath)
    mkdir(output_folder)
    new_folder = "{}/images".format(output_folder)
    mkdir(new_folder)
    output_filename = (input_filepath.split("/")[-1]).replace("pptx","ipynb")

    # load a presentation
    prs = Presentation(input_filepath)

    # Get the size
    slide_height = prs.slide_height/Inches(1)
    slide_width = prs.slide_width/Inches(1)
    if debug: print("height , width = {}, {}".format(slide_height, slide_width))

    # Initialize the cell list
    cells = [config_cell(), ]

    # Iterate through slides
    for i, slide in enumerate(prs.slides):
        if len(slide.shapes)==3:
            # Show in columns
            # Title cell
            title_content = get_markdown(slide.shapes[0], output_folder, preppend="## ")
            my_cell = create_cell(cell_type="markdown", content=title_content, slide_type="slide")
            cells.append(my_cell)
            # left_content
            if is_left_column(slide.shapes[1], slide_width) and is_right_column(slide.shapes[2], slide_width):
                left_index, right_index = 1, 2
            else:
                left_index, right_index = 2, 1
            left_markdown = get_markdown(slide.shapes[left_index], output_folder)
            right_markdown = get_markdown(slide.shapes[right_index], output_folder)
            my_text = ""
            my_text += '''<div style="float: left; width: 49%;">\n{}\n</div>'''.format(left_markdown)
            my_text += '''<div style="float: right; width: 49%;">\n{}\n</div>'''.format(right_markdown)
            my_cell = create_cell(cell_type="markdown", content=my_text, slide_type="-")
            cells.append(my_cell)
        else:
            for j, shape in enumerate(slide.shapes):
                left, width = shape.left.inches, shape.width.inches
                if debug: print("\t", i, j, "left, right, width = {}, {}, {}".format(left, left+width, width))
                top, height = shape.top.inches, shape.height.inches
                if debug: print("\t", i, j, "top, bottom, height = {}, {}, {}".format(top, top+height, height))
                if debug: print("width < slide_width? {} < {} ? {}".format(width, slide_width, width<0.5*slide_width))

                # Detect if new slide
                slide_type="slide" if j==0 else "-"
                preppend = "## " if (j==0 and get_type(shape)=="text") else ""
                my_text = get_markdown(shape, output_folder, preppend=preppend)
                my_cell = create_cell(cell_type="markdown", content=my_text, slide_type=slide_type)
                cells.append(my_cell)

    # Create notebook and save it
    ipynb = create_notebook(cells)
    save(ipynb, output_filepath)
    return

if __name__=="__main__":
    if len(sys.argv)==2:
        input_filepath=sys.argv[1]
        output_filepath = sys.argv[1].replace(".pptx",  ".ipynb")
        print("Using the default behavior")
        print("Converting {} to a jupyter notebook with RISE slides {}".format(input_filepath, output_filepath))
        ppt2rise(input_filepath, output_filepath)
    elif len(sys.argv)==3:
        input_filepath=sys.argv[1]
        output_filepath=sys.argv[2]
        print("Using the default behavior")
        print("Converting {} to a jupyter notebook with RISE slides {}".format(input_filepath, output_filepath))
        ppt2rise(input_filepath, output_filepath)
    else:
        print("Default behavior - saves the jupyter notebook in the same folder as the original pptx:")
        print("\tpython ppt2rise.py path/to/my_slides.pptx")
        print("Optional behavior - saves the jupyter notebook in the requested folder:")
        print("\tpython ppt2rise.py path/to/my_slides.pptx another/path/my_notebook.ipynb")

